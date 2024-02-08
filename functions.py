from pptx import Presentation
import re
import os
import sys

def extract_tables(pptx_file):
    """
    Extract the tables of a specified Powerpoint file (.pptx).

    Args:
        pptx_file (str): Path to the PowerPoint file.

    Returns:
        list: A list containing Table objects extracted from the pptx file.
    """
    # Check if file exists
    if not os.path.exists(pptx_file):
        print("Le fichier spécifié n'existe pas ou est mal écrit.")
        sys.exit(1)
    
    # Check file's extension
    _, extension = os.path.splitext(pptx_file)
    if extension.lower() != ".pptx":
        print("Le fichier spécifié n'est pas un fichier PowerPoint (.pptx).")
        sys.exit(1)
    
    presentation = Presentation(pptx_file)
    tables = []

    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_table:
                tables.append(shape.table)
    
    # Checks if tables were extracted
    if not tables:
        print("Aucun tableau trouvé dans le fichier PowerPoint.")
        sys.exit(1)

    return tables

def compile_regex(patterns):
    """
    Compile a list of search patterns into regular expressions.

    Args:
        patterns (list): A list of strings representing the search patterns.

    Returns:
        list: A list of objects of type `re.Pattern` representing the corresponding compiled regular expressions.
    """
    return [re.compile(pattern, re.IGNORECASE) for pattern in patterns]

def extract_matches(tables, patterns):
    """
    Searches and extracts specific information in the extracted tables.

    Args:
        tables (list): A list containing Table objects of the extracted tables from the PowerPoint file.
        patterns (list): A list of search patterns for specific information.

    Returns:
        dict: A dictionary containing the found information for each search pattern.
              The keys of the dictionary are the search patterns, and the values are the corresponding found information.
    """
    matching_dictionary = {}

    # Compile the search patterns into regular expressions
    regex_patterns = compile_regex(patterns)

    # Iterate through the extracted tables
    for table in tables:
        # Iterate through the cells of the table
        for row_index, row in enumerate(table.rows):
            for col_index, cell in enumerate(row.cells):
                # Search for information in each cell with each search pattern
                for pattern, regex_pattern in zip(patterns, regex_patterns):
                    match = regex_pattern.search(cell.text)
                    # Retrieve the content of the cell adjacent to the match
                    if match and table.cell(row_index, col_index + 1).text != "":
                        adjacent_cell = table.cell(row_index, col_index + 1)
                        matching_dictionary[pattern] = adjacent_cell.text

    return matching_dictionary


def display_information(matching_dictionary):
    """
    Display information stored in a dictionary.

    Args:
        matching_dictionary (dict): A dictionary containing information to be displayed.

    Returns:
        None
    """
    if matching_dictionary:
        for key, val in matching_dictionary.items():
            print(f"{key}: {val}")
    else:
        print("No match found")