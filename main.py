import re
from pptx import Presentation

def compile_regex(pattern):
    """
    Compile a regular expression pattern and return the compiled regex.

    :param pattern: The regular expression pattern.
    :type pattern: str

    :return: Compiled regex object.
    :rtype: re.Pattern
    """
    return re.compile(pattern, re.IGNORECASE)

def find_regex_in_paragraphs(paragraphs, regex):
    """
    Search for a regex pattern in a list of paragraphs.

    :param paragraphs: List of paragraphs to search.
    :type paragraphs: list
    :param regex: Compiled regular expression pattern.
    :type regex: re.Pattern

    :return: Matched string if found, otherwise None.
    :rtype: str or None
    """
    for paragraph in paragraphs:
        result = regex.search(paragraph.text)
        if result:
            return result.group()
    return None

def print_table_cell_content(tbl, row, col, regex_list):
    """
    Print the content of a table cell that matches any of the provided regex patterns.

    :param tbl: Table object.
    :type tbl: pptx.table.Table
    :param row: Row index of the cell.
    :type row: int
    :param col: Column index of the cell.
    :type col: int
    :param regex_list: List of compiled regex patterns.
    :type regex_list: list

    :return: Matched content if found, otherwise None.
    :rtype: str or None
    """
    cell = tbl.cell(row, col)
    paragraphs = cell.text_frame.paragraphs

    for paragraph in paragraphs:
        for regex in regex_list:
            match_result = find_regex_in_paragraphs(paragraphs, regex)
            # Check if the matched regex has content in the adjacent cell
            if match_result and tbl.cell(row, col+1).text != "":
                return f"{match_result} {tbl.cell(row, col+1).text}"
    return None

def process_table(tbl, regex_list):
    """
    Process a table, searching for regex patterns in each cell.

    :param tbl: Table object.
    :type tbl: pptx.table.Table
    :param regex_list: List of compiled regex patterns.
    :type regex_list: list
    """
    row_count = len(tbl.rows)
    col_count = len(tbl.columns)

    for r in range(row_count):
        for c in range(col_count - 1):  # Ignore the last column to avoid index out of range
            result = print_table_cell_content(tbl, r, c, regex_list)
            if result:
                print(result)

def process_slide(slide, regex_list):
    """
    Process a slide, searching for regex patterns in tables.

    :param slide: Slide object.
    :type slide: pptx.slide.Slide
    :param regex_list: List of compiled regex patterns.
    :type regex_list: list
    """
    for shape in slide.shapes:
        if shape.has_table:
            tbl = shape.table
            process_table(tbl, regex_list)

def process_presentation(prs, regex_list):
    """
    Process the entire presentation, searching for regex patterns in all slides.

    :param prs: Presentation object.
    :type prs: pptx.presentation.Presentation
    :param regex_list: List of compiled regex patterns.
    :type regex_list: list
    """
    for i, slide in enumerate(prs.slides, start=1):
        print(f"Processing Slide {i}")
        process_slide(slide, regex_list)

# Define regex patterns
pattern_FASEP = r"Montant du FASEP"
pattern_date_signature = r"Date de signature de la convention"
pattern_avis = r"Avis sur le versement intermédiaire"

# Compile regex
regex_FASEP = compile_regex(pattern_FASEP)
regex_date_signature = compile_regex(pattern_date_signature)
regex_avis = compile_regex(pattern_avis)

# Replace 'example FASEP.pptx' with the path to your PowerPoint file
pptx_file_path = 'example FASEP.pptx'
presentation = Presentation(pptx_file_path)

# Process the presentation with a list of regex patterns
process_presentation(presentation, [regex_FASEP, regex_date_signature, regex_avis])
